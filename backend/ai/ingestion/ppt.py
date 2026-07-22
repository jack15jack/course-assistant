from pptx import Presentation


def extract_ppt_text(filepath):

    prs = Presentation(filepath)

    slides = []

    for index, slide in enumerate(prs.slides):

        text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text.append(shape.text)

        slide_text = "\n".join(text)

        if slide_text.strip():

            slides.append(
                {
                    "content": slide_text,
                    "content_type": "text",
                    "metadata": {
                        "slide": index + 1,
                        "location": filepath
                    }
                }
            )

    return slides